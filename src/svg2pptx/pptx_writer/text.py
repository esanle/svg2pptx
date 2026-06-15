"""PowerPoint text box creation from SVG text elements."""

from typing import Optional

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from svg2pptx.parser.svg_parser import TextElement
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.pptx_writer.shapes import parse_hex_color


def create_text(
    shapes: SlideShapes,
    text_element: TextElement,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint text box from an SVG text element.

    Args:
        shapes: SlideShapes collection to add text box to.
        text_element: Parsed SVG text element.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created text box shape or None.
    """
    if not text_element.text:
        return None

    # Apply transform to position
    x, y = text_element.transform.apply(text_element.x, text_element.y)

    # Get font size in pixels
    font_size_px = text_element.style.font_size * scale

    # Check if this text came from foreignObject (has explicit width)
    foreign_width = getattr(text_element, "_foreign_width", None)
    foreign_height = getattr(text_element, "_foreign_height", None)

    if foreign_width and foreign_width > 0:
        # Use explicit width from foreignObject div
        estimated_width = foreign_width * scale
        estimated_height = (foreign_height or font_size_px * 1.5) * scale
    else:
        # Estimate text box dimensions for regular SVG text
        # Use a more generous width estimate to avoid text wrapping
        estimated_width = len(text_element.text) * font_size_px * 0.7 + font_size_px
        estimated_height = font_size_px * 1.4

    # Convert position to EMU
    # In SVG, text y-coordinate is the baseline position
    # We need to adjust based on the text-anchor for horizontal positioning
    text_anchor = text_element.style.text_anchor

    # Calculate left position based on text-anchor
    if text_anchor == "middle":
        # Text is centered at x
        left = offset_x + px_to_emu(x * scale) - px_to_emu(estimated_width / 2)
    elif text_anchor == "end":
        # Text ends at x
        left = offset_x + px_to_emu(x * scale) - px_to_emu(estimated_width)
    else:
        # Default: text starts at x (text-anchor="start")
        left = offset_x + px_to_emu(x * scale)

    # For y position: SVG y is the baseline, so we move up by approximately
    # the font ascent (roughly 80% of font size for most fonts)
    # For foreignObject text, y is the center position from padding-top
    baseline_offset = font_size_px * 0.85
    top = offset_y + px_to_emu(y * scale) - px_to_emu(baseline_offset)

    width = px_to_emu(estimated_width)
    height = px_to_emu(estimated_height)

    # Create text box
    text_box = shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Remove margins/padding for more accurate positioning
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    # Check for foreignObject text runs (with bold/color formatting)
    text_runs = getattr(text_element, "_text_runs", None)

    # Handle multi-line text (split by newlines)
    lines = text_element.text.split("\n")

    # Clear the first paragraph and add content
    paragraph = text_frame.paragraphs[0]

    if text_runs and len(text_runs) > 0:
        # Use text runs for rich formatting
        first_line_runs = [r for r in text_runs if "\n" not in r.text or r.text == "\n"]
        # For simplicity, combine runs for first line
        first_line_text = "".join(r.text for r in first_line_runs if r.text != "\n")
        run = paragraph.add_run()
        run.text = first_line_text
        apply_run_style(run, first_line_runs[0] if first_line_runs else None,
                       text_element, scale)

        # Add remaining lines as new paragraphs
        for line in lines[1:]:
            if line:
                p = text_frame.add_paragraph()
                r = p.add_run()
                r.text = line
                apply_run_style(r, None, text_element, scale)
    else:
        # Simple text - single run on first paragraph
        run = paragraph.add_run()
        run.text = lines[0] if lines else text_element.text
        apply_run_style(run, None, text_element, scale)

        # Add remaining lines as new paragraphs
        for line in lines[1:]:
            if line:
                p = text_frame.add_paragraph()
                r = p.add_run()
                r.text = line
                apply_run_style(r, None, text_element, scale)

    # Text anchor (horizontal alignment)
    anchor_map = {
        "start": PP_ALIGN.LEFT,
        "middle": PP_ALIGN.CENTER,
        "end": PP_ALIGN.RIGHT,
    }
    # Apply alignment to all paragraphs
    for p in text_frame.paragraphs:
        p.alignment = anchor_map.get(text_anchor, PP_ALIGN.LEFT)

    # Disable shadow on text box
    try:
        text_box.shadow.inherit = False
        if hasattr(text_box.shadow, 'visible'):
            text_box.shadow.visible = False
    except (AttributeError, NotImplementedError):
        pass

    return text_box


def apply_run_style(run, text_run, text_element, scale):
    """Apply styling to a text run."""
    font = run.font
    font.name = text_element.style.font_family
    font.size = Pt(text_element.style.font_size * scale)

    # Font weight - check both element style and run-specific bold
    is_bold = text_element.style.font_weight in ("bold", "700", "800", "900")
    if text_run and text_run.bold:
        is_bold = True
    font.bold = is_bold

    # Text color - check run-specific color override first
    if text_run and text_run.color:
        try:
            color = parse_hex_color(text_run.color)
            font.color.rgb = color
        except ValueError:
            pass
    elif text_element.style.fill != "none":
        try:
            color = parse_hex_color(text_element.style.fill)
            font.color.rgb = color
        except ValueError:
            pass

