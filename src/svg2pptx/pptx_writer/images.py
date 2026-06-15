"""PowerPoint picture shape creation from SVG image elements."""

from io import BytesIO
from typing import Optional

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes

from svg2pptx.parser.images import ImageElement
from svg2pptx.geometry.units import px_to_emu


def create_picture(
    shapes: SlideShapes,
    image_element: ImageElement,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint picture shape from an SVG image element.

    Args:
        shapes: SlideShapes collection to add picture to.
        image_element: Parsed SVG image element.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created picture shape or None if no image data.
    """
    # Need binary image data
    if not image_element.image_bytes:
        return None

    # Apply transform to position
    x, y = image_element.transform.apply(image_element.x, image_element.y)

    # Apply scale
    x_scaled = x * scale
    y_scaled = y * scale
    width_scaled = image_element.width * scale
    height_scaled = image_element.height * scale

    # Convert position and size to EMU
    left = offset_x + px_to_emu(x_scaled)
    top = offset_y + px_to_emu(y_scaled)
    width = px_to_emu(width_scaled)
    height = px_to_emu(height_scaled)

    # Skip if dimensions are invalid
    if width <= 0 or height <= 0:
        return None

    # Create BytesIO stream for the image data
    image_stream = BytesIO(image_element.image_bytes)

    # Determine content type for python-pptx
    # python-pptx infers type from file signature, but we can help it
    mime_type = image_element.mime_type.lower()

    try:
        # add_picture accepts BytesIO
        picture = shapes.add_picture(
            image_stream,
            left,
            top,
            width,
            height,
        )
        return picture
    except Exception as e:
        # Log warning but don't crash
        print(f"Warning: Failed to add image: {e}")
        return None